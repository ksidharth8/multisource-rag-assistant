from pathlib import Path
import fitz
from docx import Document
from pptx import Presentation


SUPPORTED_DOCUMENT_EXTENSIONS = {".pdf", ".txt", ".md", ".docx", ".pptx"}


def load_pdf(path: str) -> str:
    parts: list[str] = []
    with fitz.open(path) as doc:
        for page_number, page in enumerate(doc, start=1):
            text = page.get_text("text")
            if text.strip():
                parts.append(f"\n[Page {page_number}]\n{text}")
    return "\n".join(parts)


def load_txt(path: str) -> str:
    return Path(path).read_text(encoding="utf-8", errors="ignore")


def load_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def load_pptx(path: str) -> str:
    prs = Presentation(path)
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"\n[Slide {idx}]\n" + "\n".join(texts))
    return "\n".join(slides)


def load_document(path: str, filename: str | None = None) -> str:
    ext = Path(filename or path).suffix.lower()
    if ext == ".pdf":
        return load_pdf(path)
    if ext in {".txt", ".md"}:
        return load_txt(path)
    if ext == ".docx":
        return load_docx(path)
    if ext == ".pptx":
        return load_pptx(path)
    raise ValueError(f"Unsupported document type: {ext}")


def load_document_text(path: str, filename: str | None = None) -> str:
    return load_document(path, filename)